#!/usr/bin/env python3
"""
Generate Sampler-Crow-Build-Guide.pptx from the current hardware design.
Run from the project root:  python3 scripts/build_guide_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- Theme ----------
BG = RGBColor(0x1A, 0x1A, 0x2E)
PANEL = RGBColor(0x22, 0x22, 0x44)
ACCENT = RGBColor(0x00, 0xD4, 0xFF)
ACCENT_DIM = RGBColor(0x00, 0x88, 0xAA)
WHITE = RGBColor(0xE0, 0xE0, 0xE0)
DIM = RGBColor(0x88, 0x88, 0x88)
GREEN = RGBColor(0x44, 0xFF, 0x44)
ORANGE = RGBColor(0xFF, 0x88, 0x44)
RED = RGBColor(0xFF, 0x44, 0x44)
YELLOW = RGBColor(0xFF, 0xFF, 0x44)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def blank_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    return s

def add_text(slide, x, y, w, h, text, size=14, color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, font='SF Mono'):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)

    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        f = r.font
        f.name = font
        f.size = Pt(size)
        f.color.rgb = color
        f.bold = bold
    return box

def add_title(slide, text, subtitle=None):
    add_text(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6),
             text, size=28, color=ACCENT, bold=True, font='SF Pro Display')
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.9), Inches(12.3), Inches(0.4),
                 subtitle, size=14, color=DIM, font='SF Pro Display')
    # Accent underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.35),
                                   Inches(0.8), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

def add_panel(slide, x, y, w, h, color=PANEL):
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    p.fill.solid()
    p.fill.fore_color.rgb = color
    p.line.color.rgb = ACCENT_DIM
    p.line.width = Pt(0.5)
    return p

def add_hyperlink(slide, x, y, w, h, text, url, size=11):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = 'SF Mono'
    f.size = Pt(size)
    f.color.rgb = ACCENT
    f.underline = True
    r.hyperlink.address = url
    return box

# ==========================================================
# SLIDE 1 — Cover
# ==========================================================
s = blank_slide()
# Big title
add_text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.2),
         'Sampler-Crow', size=72, color=ACCENT, bold=True,
         align=PP_ALIGN.CENTER, font='SF Pro Display')
add_text(s, Inches(0.5), Inches(2.8), Inches(12.3), Inches(0.6),
         'Portable Music Workstation — Build Guide v3',
         size=24, color=WHITE, align=PP_ALIGN.CENTER, font='SF Pro Display')
add_text(s, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.5),
         'Teensy 4.1  +  Audio Shield  +  CrowPanel 5"  +  Novation Launchpad',
         size=16, color=DIM, align=PP_ALIGN.CENTER, font='SF Pro Display')

# Feature boxes
features = [
    ('8-TRACK\nSEQUENCER', GREEN),
    ('USB MIDI\nHOST', ORANGE),
    ('TRS LINE\nIN/OUT', ACCENT),
    ('SAMPLE\nPLAYBACK', YELLOW),
]
box_w = Inches(2.6)
box_h = Inches(1.3)
gap = Inches(0.2)
total_w = box_w * len(features) + gap * (len(features) - 1)
start_x = (prs.slide_width - total_w) / 2
for i, (label, color) in enumerate(features):
    bx = start_x + i * (box_w + gap)
    by = Inches(4.5)
    panel = add_panel(s, bx, by, box_w, box_h)
    panel.line.color.rgb = color
    panel.line.width = Pt(2)
    add_text(s, bx, by + Inches(0.3), box_w, Inches(0.8),
             label, size=16, color=color, bold=True,
             align=PP_ALIGN.CENTER, font='SF Pro Display')

add_text(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.5),
         'Updated: 2026-04-11  •  github.com/nohumans-com/Sampler-Crow',
         size=11, color=DIM, align=PP_ALIGN.CENTER, font='SF Mono')

# ==========================================================
# SLIDE 2 — What you'll build
# ==========================================================
s = blank_slide()
add_title(s, 'What you will build',
          'A standalone drum machine + sampler + synth that runs without a computer')

bullets = [
    '8 independent tracks (kick / snare / hats / clap / bass / lead / pluck / fx)',
    'Step sequencer with 8 steps per track, BPM clock, pattern save/load',
    'Real Novation Launchpad Mini MK3 as the step editor and clip launcher',
    'CrowPanel 5" touchscreen as the visual UI (mixer, piano roll, samples)',
    'Analog stereo TRS line in, line out, and headphone out',
    '5-pin TRS MIDI in and MIDI out',
    'Companion macOS host app for development, monitoring, sample loading',
]
y = Inches(1.8)
for b in bullets:
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y + Inches(0.08),
                             Inches(0.12), Inches(0.12))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT
    dot.line.fill.background()
    add_text(s, Inches(1.0), y, Inches(11.5), Inches(0.4), b,
             size=15, color=WHITE, font='SF Pro Display')
    y += Inches(0.55)

# ==========================================================
# SLIDE 3 — Architecture diagram
# ==========================================================
s = blank_slide()
add_title(s, 'System architecture',
          'Two processors + Audio Shield + powered USB hub')

# CrowPanel at top
add_panel(s, Inches(4.5), Inches(1.8), Inches(4.3), Inches(0.9))
add_text(s, Inches(4.5), Inches(1.85), Inches(4.3), Inches(0.4),
         'CrowPanel 5.0" HMI', size=14, color=ACCENT, bold=True,
         align=PP_ALIGN.CENTER, font='SF Pro Display')
add_text(s, Inches(4.5), Inches(2.2), Inches(4.3), Inches(0.4),
         'ESP32-S3-WROOM-1  |  800×480 touch  |  LVGL UI',
         size=10, color=DIM, align=PP_ALIGN.CENTER, font='SF Mono')

# Teensy center
add_panel(s, Inches(4.0), Inches(3.3), Inches(5.3), Inches(2.4))
add_text(s, Inches(4.0), Inches(3.4), Inches(5.3), Inches(0.4),
         'Teensy 4.1  +  Audio Shield', size=16, color=ACCENT, bold=True,
         align=PP_ALIGN.CENTER, font='SF Pro Display')
add_text(s, Inches(4.0), Inches(3.75), Inches(5.3), Inches(0.3),
         'ARM Cortex-M7 @ 600 MHz  •  SGTL5000 codec',
         size=10, color=DIM, align=PP_ALIGN.CENTER, font='SF Mono')
add_text(s, Inches(4.15), Inches(4.1), Inches(5.0), Inches(1.5),
         'Audio engine  •  Sequencer  •  Sample playback\n'
         'USB Host (MIDI)  •  USB Device (Audio+MIDI+Serial)\n'
         'MIDI TRS UART  •  SD card  •  4 pots, 4 buttons',
         size=11, color=WHITE, align=PP_ALIGN.CENTER, font='SF Mono')

# Peripherals around Teensy
peripheral_positions = [
    ('LAUNCHPAD\nMini MK3',     Inches(0.5),  Inches(4.0), ORANGE),
    ('USB HUB\n(powered)',       Inches(0.5),  Inches(5.2), DIM),
    ('MIDI TRS\nin/out',         Inches(10.0), Inches(4.0), YELLOW),
    ('LINE TRS\nin/out/hp',      Inches(10.0), Inches(5.2), GREEN),
    ('µSD\nsamples',             Inches(4.0),  Inches(6.0), ACCENT),
    ('macOS\nhost app',          Inches(8.5),  Inches(6.0), RED),
]
for label, x, y, color in peripheral_positions:
    p = add_panel(s, x, y, Inches(1.8), Inches(0.9))
    p.line.color.rgb = color
    add_text(s, x, y + Inches(0.1), Inches(1.8), Inches(0.7),
             label, size=10, color=color, bold=True,
             align=PP_ALIGN.CENTER, font='SF Pro Display')

# ==========================================================
# SLIDE 4 — Bill of Materials (part 1)
# ==========================================================
s = blank_slide()
add_title(s, 'Bill of Materials (1/2)',
          'Core components + audio/MIDI I/O  •  all-in total ~$245')

# Table headers
header_y = Inches(1.8)
cols = [('Part', Inches(0.5), Inches(3.5)),
        ('Product', Inches(4.0), Inches(5.5)),
        ('Price', Inches(9.5), Inches(1.0)),
        ('Link', Inches(10.5), Inches(2.3))]
for title, x, w in cols:
    add_text(s, x, header_y, w, Inches(0.3), title, size=12,
             color=ACCENT, bold=True, font='SF Pro Display')

# Divider
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.15),
                          Inches(12.3), Inches(0.02))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT_DIM; line.line.fill.background()

bom_rows = [
    ('Teensy 4.1 (with pins)',    'PJRC Teensy 4.1 ARM Cortex-M7 600 MHz', '$49',
     'https://www.pjrc.com/store/teensy41_pins.html'),
    ('Audio Shield Rev D',        'Teensy Audio Shield (SGTL5000)',        '$14',
     'https://www.pjrc.com/store/teensy3_audio.html'),
    ('CrowPanel 5.0" HMI',        'Elecrow ESP32-S3 800x480 cap touch',    '$60',
     'https://www.elecrow.com/esp32-display-5-inch-hmi-display-rgb-tft-lcd-touch-screen-support-lvgl.html'),
    ('Novation Launchpad Mini 3', 'Novation LAUNCHPADMINIMK3',             '$109',
     'https://www.amazon.com/dp/B086ZKJFNC'),
    ('ubld.it MIDI Breakout MV',  'Opto-isolated MIDI I/O, 3.3/5 V',       '$25',
     'https://www.amazon.com/dp/B0BYMC926Z'),
    ('3.5 mm TRS jacks (5-pack)', 'Panel mount stereo jacks',              '$7',
     'https://www.amazon.com/dp/B07JNC4P7Y'),
    ('MicroSD 32 GB',             'SanDisk Ultra 32 GB',                   '$8',
     'https://www.amazon.com/dp/B08GY9NYRM'),
]
y = Inches(2.3)
for part, product, price, url in bom_rows:
    add_text(s, Inches(0.5), y, Inches(3.5), Inches(0.3),
             part, size=11, color=WHITE, font='SF Pro Display', bold=True)
    add_text(s, Inches(4.0), y, Inches(5.5), Inches(0.3),
             product, size=10, color=DIM, font='SF Mono')
    add_text(s, Inches(9.5), y, Inches(1.0), Inches(0.3),
             price, size=11, color=GREEN, bold=True, font='SF Mono')
    add_hyperlink(s, Inches(10.5), y, Inches(2.3), Inches(0.3),
                  'View product →', url, size=10)
    y += Inches(0.55)

# ==========================================================
# SLIDE 5 — Bill of Materials (part 2) — Power
# ==========================================================
s = blank_slide()
add_title(s, 'Bill of Materials (2/2)',
          'USB host + power + controls  •  pick ONE power option')

header_y = Inches(1.6)
for title, x, w in cols:
    add_text(s, x, header_y, w, Inches(0.3), title, size=12,
             color=ACCENT, bold=True, font='SF Pro Display')
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.95),
                          Inches(12.3), Inches(0.02))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT_DIM; line.line.fill.background()

bom2_rows = [
    ('SparkFun USB Host Cable',    '5-pin to USB-A female (Teensy 4.1)',  '$4',
     'https://www.sparkfun.com/usb-host-cable-for-teensy-4-1-and-teensy-3-6.html'),
    ('Powered USB hub',            'Anker 4-Port USB 2.0 Hub with AC adapter', '$15',
     'https://www.amazon.com/dp/B00DQFGH80'),
    ('USB-C panel mount (2x)',     'USB-C male-to-female pigtails',        '$8',
     'https://www.amazon.com/dp/B08HS6X44P'),
    ('B10K potentiometers',        '×4 with knobs (20-pack)',              '$9',
     'https://www.amazon.com/dp/B06WWQP12J'),
    ('Tactile buttons',            '12 mm momentary ×4 (25-pack)',         '$7',
     'https://www.amazon.com/dp/B0798HZ8WB'),
    ('Perfboard + hookup wire',    'Proto PCB + 22 AWG solid core',        '$18',
     'https://www.amazon.com/dp/B07W83VJGV'),
]
y = Inches(2.1)
for part, product, price, url in bom2_rows:
    add_text(s, Inches(0.5), y, Inches(3.5), Inches(0.3),
             part, size=11, color=WHITE, font='SF Pro Display', bold=True)
    add_text(s, Inches(4.0), y, Inches(5.5), Inches(0.3),
             product, size=10, color=DIM, font='SF Mono')
    add_text(s, Inches(9.5), y, Inches(1.0), Inches(0.3),
             price, size=11, color=GREEN, bold=True, font='SF Mono')
    add_hyperlink(s, Inches(10.5), y, Inches(2.3), Inches(0.3),
                  'View product →', url, size=10)
    y += Inches(0.45)

# Power options header
add_text(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.4),
         'POWER — pick ONE of these four options',
         size=14, color=ORANGE, bold=True, font='SF Pro Display')
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.35),
                          Inches(12.3), Inches(0.02))
line.fill.solid(); line.fill.fore_color.rgb = ORANGE; line.line.fill.background()

power_rows = [
    ('A: Wall wart 5V/3A',   'Mean Well GSM18A05-P1J + barrel jack + perfboard',
     '$19', 'https://www.digikey.com/en/products/detail/mean-well-usa-inc/GSM18A05-P1J/7703632'),
    ('B: Internal AC-DC',    'Mean Well IRM-20-5 (20W, 5V/4A, enclosed) + C8 inlet',
     '$21', 'https://www.digikey.com/en/products/detail/mean-well-usa-inc/IRM-20-5/7704667'),
    ('C: USB-C PD input',    'Adafruit HUSB238 breakout (any PD charger)',
     '$14', 'https://www.adafruit.com/product/5807'),
    ('D: Buck regulator',    'Pololu D36V28F5 (5V/3.2A) + 9-12V brick',
     '$25', 'https://www.pololu.com/product/3782'),
]
y = Inches(5.55)
for part, product, price, url in power_rows:
    add_text(s, Inches(0.5), y, Inches(3.5), Inches(0.3),
             part, size=11, color=WHITE, font='SF Pro Display', bold=True)
    add_text(s, Inches(4.0), y, Inches(5.5), Inches(0.3),
             product, size=10, color=DIM, font='SF Mono')
    add_text(s, Inches(9.5), y, Inches(1.0), Inches(0.3),
             price, size=11, color=GREEN, bold=True, font='SF Mono')
    add_hyperlink(s, Inches(10.5), y, Inches(2.3), Inches(0.3),
                  'View product →', url, size=10)
    y += Inches(0.4)

# ==========================================================
# SLIDE 6 — Power strategy detail
# ==========================================================
s = blank_slide()
add_title(s, 'Power strategy',
          'Current budget  •  distribution  •  critical wiring rules')

# Current budget table on left
add_text(s, Inches(0.5), Inches(1.8), Inches(5.0), Inches(0.3),
         'Current budget @ 5 V', size=14, color=ACCENT, bold=True,
         font='SF Pro Display')
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.1),
                          Inches(5.5), Inches(0.02))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT_DIM; line.line.fill.background()

loads = [
    ('Teensy 4.1 + Audio Shield',      '120 mA',  '180 mA'),
    ('CrowPanel 5" (backlight)',       '250 mA',  '400 mA'),
    ('Launchpad Mini MK3',             '200 mA',  '500 mA'),
    ('Powered USB hub + expansion',    '30 mA',   '550 mA'),
    ('Pots, buttons, misc',            '10 mA',   '10 mA'),
]
y = Inches(2.3)
add_text(s, Inches(0.5), y, Inches(3.5), Inches(0.3),
         'Load', size=10, color=DIM, bold=True, font='SF Mono')
add_text(s, Inches(4.0), y, Inches(1.0), Inches(0.3),
         'Typical', size=10, color=DIM, bold=True, font='SF Mono')
add_text(s, Inches(5.1), y, Inches(1.0), Inches(0.3),
         'Peak', size=10, color=DIM, bold=True, font='SF Mono')
y += Inches(0.35)
for label, typ, peak in loads:
    add_text(s, Inches(0.5), y, Inches(3.5), Inches(0.3),
             label, size=11, color=WHITE, font='SF Pro Display')
    add_text(s, Inches(4.0), y, Inches(1.0), Inches(0.3),
             typ, size=11, color=YELLOW, font='SF Mono')
    add_text(s, Inches(5.1), y, Inches(1.0), Inches(0.3),
             peak, size=11, color=ORANGE, font='SF Mono')
    y += Inches(0.35)
# Total row
line2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(5.5), Inches(0.02))
line2.fill.solid(); line2.fill.fore_color.rgb = ACCENT_DIM; line2.line.fill.background()
y += Inches(0.1)
add_text(s, Inches(0.5), y, Inches(3.5), Inches(0.3),
         'TOTAL', size=12, color=WHITE, bold=True, font='SF Pro Display')
add_text(s, Inches(4.0), y, Inches(1.0), Inches(0.3),
         '~610 mA', size=12, color=GREEN, bold=True, font='SF Mono')
add_text(s, Inches(5.1), y, Inches(1.0), Inches(0.3),
         '~1.64 A', size=12, color=RED, bold=True, font='SF Mono')

y += Inches(0.6)
add_text(s, Inches(0.5), y, Inches(5.5), Inches(0.3),
         'Recommend 5 V / 3 A supply (comfortable headroom)',
         size=11, color=GREEN, bold=True, font='SF Mono')

# Critical rules on right
add_text(s, Inches(6.5), Inches(1.8), Inches(6.0), Inches(0.3),
         'Critical wiring rules', size=14, color=ORANGE, bold=True,
         font='SF Pro Display')
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(2.1),
                          Inches(6.3), Inches(0.02))
line.fill.solid(); line.fill.fore_color.rgb = ORANGE; line.line.fill.background()

rules = [
    ('1.', 'Cut the VUSB→VIN pad underneath the Teensy 4.1 before\n    powering from VIN. Standard PJRC procedure, one scalpel cut.'),
    ('2.', 'DO NOT back-feed the USB-C COMPUTER port from your\n    5 V rail. It is a device port. Leave VBUS for enumeration only.'),
    ('3.', 'The USB hub must be POWERED from the system 5 V rail\n    via its own polyfuse, NOT from the Teensy USB host 5V pad.'),
    ('4.', 'Add a 1N5819 Schottky diode on the barrel jack input for\n    reverse-polarity protection. Cheap insurance.'),
    ('5.', 'Add a 1000 µF bulk cap near the distribution block to\n    absorb the Launchpad LED inrush spike.'),
]
y = Inches(2.3)
for num, rule in rules:
    add_text(s, Inches(6.5), y, Inches(0.3), Inches(0.4),
             num, size=12, color=ORANGE, bold=True, font='SF Mono')
    add_text(s, Inches(6.9), y, Inches(6.0), Inches(0.8),
             rule, size=11, color=WHITE, font='SF Pro Display')
    y += Inches(0.9)

# ==========================================================
# SLIDE 7 — Panel layout
# ==========================================================
s = blank_slide()
add_title(s, 'Panel layout',
          'Front = display + controls  •  Rear = all I/O')

# Front panel
add_text(s, Inches(0.5), Inches(1.8), Inches(12.3), Inches(0.3),
         'FRONT PANEL', size=12, color=ACCENT, bold=True, font='SF Pro Display')
front = add_panel(s, Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.8))
# display rect
d = add_panel(s, Inches(0.9), Inches(2.5), Inches(5.5), Inches(1.2), PANEL)
d.line.color.rgb = ACCENT; d.line.width = Pt(2)
add_text(s, Inches(0.9), Inches(2.8), Inches(5.5), Inches(0.6),
         '5" Touchscreen  800×480',
         size=14, color=ACCENT, bold=True,
         align=PP_ALIGN.CENTER, font='SF Pro Display')
# pots
for i in range(4):
    px = Inches(6.8 + i * 0.8)
    p = s.shapes.add_shape(MSO_SHAPE.OVAL, px, Inches(2.6), Inches(0.6), Inches(0.6))
    p.fill.solid(); p.fill.fore_color.rgb = PANEL
    p.line.color.rgb = ORANGE; p.line.width = Pt(1.5)
    add_text(s, px, Inches(3.25), Inches(0.6), Inches(0.3),
             f'P{i+1}', size=10, color=DIM,
             align=PP_ALIGN.CENTER, font='SF Mono')
# buttons
for i in range(4):
    bx = Inches(10.2 + i * 0.55)
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, Inches(2.7),
                           Inches(0.4), Inches(0.4))
    b.fill.solid(); b.fill.fore_color.rgb = PANEL
    b.line.color.rgb = GREEN; b.line.width = Pt(1.5)
    add_text(s, bx, Inches(3.15), Inches(0.4), Inches(0.3),
             f'B{i+1}', size=9, color=DIM,
             align=PP_ALIGN.CENTER, font='SF Mono')

# Rear panel
add_text(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.3),
         'REAR PANEL', size=12, color=ACCENT, bold=True, font='SF Pro Display')
rear = add_panel(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(1.8))

ports = [
    ('5V DC', 'IN',         RED),
    ('USB-C', 'COMPUTER',   ACCENT),
    ('USB-C', 'CTRLR',      ORANGE),
    ('LINE',  'IN',         GREEN),
    ('LINE',  'OUT',        GREEN),
    ('HEAD',  'PHONE',      GREEN),
    ('MIDI',  'OUT',        YELLOW),
    ('MIDI',  'IN',         YELLOW),
]
port_w = Inches(1.4)
port_gap = Inches(0.1)
total_w = port_w * len(ports) + port_gap * (len(ports) - 1)
start_x = Inches(0.5) + (Inches(12.3) - total_w) / 2
for i, (label1, label2, color) in enumerate(ports):
    px = start_x + i * (port_w + port_gap)
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, px + Inches(0.35), Inches(5.0),
                                Inches(0.7), Inches(0.7))
    circle.fill.solid(); circle.fill.fore_color.rgb = PANEL
    circle.line.color.rgb = color; circle.line.width = Pt(2)
    add_text(s, px, Inches(5.8), port_w, Inches(0.3),
             label1, size=10, color=color, bold=True,
             align=PP_ALIGN.CENTER, font='SF Mono')
    add_text(s, px, Inches(6.1), port_w, Inches(0.3),
             label2, size=9, color=DIM,
             align=PP_ALIGN.CENTER, font='SF Mono')

add_text(s, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
         'Note: USB-C COMPUTER is a DEVICE port (host your laptop sees). '
         'USB-C CTRLR is for the Launchpad via the internal powered hub.',
         size=10, color=DIM, align=PP_ALIGN.CENTER, font='SF Pro Display')

# ==========================================================
# SLIDE 8 — Wiring: Audio Shield + controls
# ==========================================================
s = blank_slide()
add_title(s, 'Step 1: Audio Shield + controls',
          'Stack the Audio Shield, wire 4 pots and 4 buttons, verify tone output')

add_text(s, Inches(0.5), Inches(1.8), Inches(12.3), Inches(0.4),
         '1. Stack the Audio Shield onto the Teensy 4.1 (no soldering — just press).',
         size=13, color=WHITE, font='SF Pro Display')
add_text(s, Inches(0.5), Inches(2.25), Inches(12.3), Inches(0.4),
         '2. Wire each pot: outer legs to GND and 3.3 V, wiper to analog input.',
         size=13, color=WHITE, font='SF Pro Display')
add_text(s, Inches(0.5), Inches(2.7), Inches(12.3), Inches(0.4),
         '3. Wire each button: pin to GND (uses INPUT_PULLUP in software).',
         size=13, color=WHITE, font='SF Pro Display')
add_text(s, Inches(0.5), Inches(3.15), Inches(12.3), Inches(0.4),
         '4. Upload firmware from firmware/ directory, hear 440 Hz test tone.',
         size=13, color=WHITE, font='SF Pro Display')

# Pin table
add_text(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.3),
         'Pin map', size=13, color=ACCENT, bold=True, font='SF Pro Display')
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.3),
                          Inches(12.3), Inches(0.02))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT_DIM; line.line.fill.background()

pin_rows = [
    ('Pin A2 (16)', 'Pot 1 wiper',          'Cutoff'),
    ('Pin A3 (17)', 'Pot 2 wiper',          'Resonance'),
    ('Pin A4 (24)', 'Pot 3 wiper',          'Attack / Decay'),
    ('Pin A5 (25)', 'Pot 4 wiper',          'Volume / Mix'),
    ('Pin D2',      'Button 1 → GND',       'Play / Trigger'),
    ('Pin D6',      'Button 2 → GND',       'Record'),
    ('Pin D9',      'Button 3 → GND',       'Mode / Preset'),
    ('Pin D22',     'Button 4 → GND',       'Shift / Function'),
]
y = Inches(4.45)
for pin, wire, func in pin_rows:
    add_text(s, Inches(0.5), y, Inches(2.5), Inches(0.3),
             pin, size=11, color=ACCENT, font='SF Mono', bold=True)
    add_text(s, Inches(3.2), y, Inches(4.5), Inches(0.3),
             wire, size=11, color=WHITE, font='SF Mono')
    add_text(s, Inches(8.0), y, Inches(4.8), Inches(0.3),
             func, size=11, color=DIM, font='SF Mono')
    y += Inches(0.3)

# ==========================================================
# SLIDE 9 — Wiring: USB host + MIDI
# ==========================================================
s = blank_slide()
add_title(s, 'Step 2: USB host + MIDI',
          'Solder the 5-pin USB host header, wire the MIDI breakout board')

add_text(s, Inches(0.5), Inches(1.8), Inches(12.3), Inches(0.4),
         'USB host pads (bottom of Teensy 4.1)',
         size=13, color=ACCENT, bold=True, font='SF Pro Display')
add_text(s, Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.2),
         '• Solder 5-pin male header to the bottom pads labeled "USB HOST"\n'
         '• Connect SparkFun USB Host Cable or direct-wire to powered USB hub\n'
         '• IMPORTANT: use only the D+, D-, GND pads — leave the 5V pad disconnected\n'
         '• Wire hub VBUS IN directly from system 5 V rail through a 1 A polyfuse',
         size=12, color=WHITE, font='SF Pro Display')

add_text(s, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.4),
         'ubld.it MIDI breakout (TRS Type A)',
         size=13, color=ACCENT, bold=True, font='SF Pro Display')
add_text(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(1.6),
         '• Set the board voltage switch to 3.3 V\n'
         '• ubld.it TX ← Teensy Pin 14 (Serial3 TX)     — MIDI OUT data\n'
         '• ubld.it RX → Teensy Pin 15 (Serial3 RX)     — MIDI IN data\n'
         '• ubld.it VCC → Teensy 3.3 V\n'
         '• ubld.it GND → Teensy GND\n'
         '• Panel-mount the two 3.5 mm TRS jacks to the board\'s MIDI IN / OUT',
         size=11, color=WHITE, font='SF Mono')

add_text(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.4),
         '⚠ TRS Type A standard: Tip = current source, Ring = current sink, Sleeve = GND',
         size=11, color=ORANGE, font='SF Pro Display')

# ==========================================================
# SLIDE 10 — Wiring: Power distribution
# ==========================================================
s = blank_slide()
add_title(s, 'Step 3: Power distribution',
          'Wire up the barrel jack → distribution board → everything')

add_text(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(0.4),
         'Build the tiny distribution board on a piece of perfboard',
         size=13, color=ACCENT, bold=True, font='SF Pro Display')

steps = [
    ('1.', 'Cut the VUSB→VIN pad on the underside of the Teensy 4.1 with a sharp hobby knife.'),
    ('2.', 'Panel-mount the 2.1 × 5.5 mm barrel jack on the rear panel.'),
    ('3.', 'On perfboard: barrel jack (+) → 1N5819 Schottky diode → node X.'),
    ('4.', 'From node X: 1000 µF / 16 V electrolytic cap (+ to X, – to GND).'),
    ('5.', 'From node X: 1.5 A polyfuse → 5 V MAIN rail screw terminal.'),
    ('6.', 'From node X: 1 A polyfuse → 5 V HUB rail screw terminal.'),
    ('7.', 'Wire 5 V MAIN to Teensy VIN (5 V pin) and CrowPanel 5 V in.'),
    ('8.', 'Wire 5 V HUB to the powered USB hub\'s VBUS input.'),
    ('9.', 'Wire barrel jack (−) and everything else GND to a common ground bus.'),
    ('10.', 'Verify with multimeter: 5 V ± 0.1 V at Teensy VIN under load.'),
]
y = Inches(2.15)
for num, step in steps:
    add_text(s, Inches(0.5), y, Inches(0.4), Inches(0.3),
             num, size=11, color=ORANGE, bold=True, font='SF Mono')
    add_text(s, Inches(1.0), y, Inches(11.8), Inches(0.3),
             step, size=11, color=WHITE, font='SF Pro Display')
    y += Inches(0.4)

add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
         '⚠ Leave the USB-C COMPUTER port\'s VBUS disconnected. Do NOT back-feed.',
         size=11, color=RED, font='SF Pro Display', bold=True)

# ==========================================================
# SLIDE 11 — Power strategy decision tree
# ==========================================================
s = blank_slide()
add_title(s, 'Which power option should you choose?',
          'Four real products, pick based on your budget and aesthetic')

boxes = [
    ('A', 'Wall wart + barrel jack',
     'Cheapest • simplest • works',
     'Mean Well GSM18A05-P1J\n5 V / 3 A external brick\n+ barrel jack + small perfboard',
     '$19', 'Recommended for prototype',
     GREEN),
    ('B', 'Internal AC-DC module',
     'Professional • self-contained',
     'Mean Well IRM-20-5\n5 V / 4 A encapsulated\n+ C8 mains inlet',
     '$21', 'Recommended for final build',
     ACCENT),
    ('C', 'USB-C PD input',
     'Modern • phone charger works',
     'Adafruit HUSB238 breakout\nnegotiates 5 V / 3 A from\nany USB-C PD charger',
     '$14', 'If you want a single USB cable',
     ORANGE),
    ('D', 'Buck regulator',
     'Maximum headroom',
     'Pololu D36V28F5\n5 V / 3.2 A switcher\n+ 9-12 V wall brick',
     '$25', 'If you plan to add a power amp',
     YELLOW),
]
box_w = Inches(3.0)
box_h = Inches(4.5)
gap = Inches(0.1)
total = box_w * 4 + gap * 3
start_x = (prs.slide_width - total) / 2
by = Inches(1.9)
for i, (letter, title, tagline, detail, price, rec, color) in enumerate(boxes):
    bx = start_x + i * (box_w + gap)
    p = add_panel(s, bx, by, box_w, box_h)
    p.line.color.rgb = color; p.line.width = Pt(2)
    # Letter badge
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL, bx + Inches(0.2), by + Inches(0.2),
                               Inches(0.6), Inches(0.6))
    badge.fill.solid(); badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    add_text(s, bx + Inches(0.2), by + Inches(0.3), Inches(0.6), Inches(0.4),
             letter, size=16, color=BG, bold=True,
             align=PP_ALIGN.CENTER, font='SF Pro Display')
    # Title
    add_text(s, bx + Inches(0.95), by + Inches(0.25), box_w - Inches(1.1), Inches(0.35),
             title, size=13, color=color, bold=True, font='SF Pro Display')
    # Tagline
    add_text(s, bx + Inches(0.2), by + Inches(0.95), box_w - Inches(0.4), Inches(0.35),
             tagline, size=11, color=DIM, font='SF Pro Display')
    # Detail
    add_text(s, bx + Inches(0.2), by + Inches(1.5), box_w - Inches(0.4), Inches(1.6),
             detail, size=11, color=WHITE, font='SF Mono')
    # Price
    add_text(s, bx + Inches(0.2), by + Inches(3.1), box_w - Inches(0.4), Inches(0.4),
             price, size=24, color=color, bold=True,
             align=PP_ALIGN.CENTER, font='SF Pro Display')
    # Recommendation
    add_text(s, bx + Inches(0.2), by + Inches(3.75), box_w - Inches(0.4), Inches(0.6),
             rec, size=10, color=DIM,
             align=PP_ALIGN.CENTER, font='SF Pro Display')

# ==========================================================
# SLIDE 12 — Firmware & app setup
# ==========================================================
s = blank_slide()
add_title(s, 'Step 4: Firmware & host app',
          'Flash the Teensy, build the macOS app, start making beats')

add_text(s, Inches(0.5), Inches(1.8), Inches(6.0), Inches(0.4),
         'Firmware (Teensy 4.1)', size=14, color=ACCENT, bold=True,
         font='SF Pro Display')
add_text(s, Inches(0.5), Inches(2.2), Inches(6.0), Inches(0.3),
         'Requires PlatformIO with Python 3.13 (not 3.14)',
         size=10, color=DIM, font='SF Pro Display')

fw_panel = add_panel(s, Inches(0.5), Inches(2.6), Inches(6.0), Inches(3.5), BG)
fw_panel.line.color.rgb = ACCENT_DIM
add_text(s, Inches(0.7), Inches(2.8), Inches(5.8), Inches(3.3),
         '$ cd ~/Sampler-Crow\n'
         '$ pio run\n\n'
         '# Put Teensy in bootloader mode\n'
         '# (press the button on the board)\n\n'
         '$ teensy_loader_cli --mcu=TEENSY41 \\\n'
         '    -w -v .pio/build/teensy41/firmware.hex\n\n'
         '# Or just:\n'
         '$ pio run -t upload',
         size=11, color=GREEN, font='SF Mono')

add_text(s, Inches(7.0), Inches(1.8), Inches(6.0), Inches(0.4),
         'macOS host app', size=14, color=ACCENT, bold=True,
         font='SF Pro Display')
add_text(s, Inches(7.0), Inches(2.2), Inches(6.0), Inches(0.3),
         'SwiftUI • Swift 6 • macOS 15+',
         size=10, color=DIM, font='SF Pro Display')

app_panel = add_panel(s, Inches(7.0), Inches(2.6), Inches(5.8), Inches(3.5), BG)
app_panel.line.color.rgb = ACCENT_DIM
add_text(s, Inches(7.2), Inches(2.8), Inches(5.6), Inches(3.3),
         '$ cd SamplerCrowApp\n'
         '$ swift build -c release\n\n'
         '# Deploy to /Applications\n'
         '$ cp .build/release/SamplerCrowApp \\\n'
         '    "/Applications/Sampler Crow.app/\\\n'
         '      Contents/MacOS/Sampler Crow"\n'
         '$ xattr -cr "/Applications/Sampler Crow.app"\n'
         '$ codesign --force --sign - \\\n'
         '    "/Applications/Sampler Crow.app/\\\n'
         '      Contents/MacOS/Sampler Crow"\n'
         '$ open "/Applications/Sampler Crow.app"',
         size=11, color=GREEN, font='SF Mono')

add_text(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
         'Plug in a Novation Launchpad Mini MK3 → pick your audio interface in the Audio tab → press Play.',
         size=13, color=WHITE, align=PP_ALIGN.CENTER, font='SF Pro Display')

# ==========================================================
# SLIDE 13 — The USB audio host reality check
# ==========================================================
s = blank_slide()
add_title(s, 'Why we are NOT hosting a USB audio interface',
          'The honest technical answer — read this before asking again')

add_text(s, Inches(0.5), Inches(1.8), Inches(12.3), Inches(0.4),
         'What you might think you want',
         size=13, color=ORANGE, bold=True, font='SF Pro Display')
add_text(s, Inches(0.5), Inches(2.15), Inches(12.3), Inches(0.6),
         '"Plug a Universal Audio Volt 476 into the Sampler-Crow\'s USB host port, '
         'get 8 channels of high-end preamp inputs routed straight into the sampler engine."',
         size=12, color=WHITE, font='SF Pro Display')

add_text(s, Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.4),
         'Why it does not work on the Teensy 4.1',
         size=13, color=RED, bold=True, font='SF Pro Display')
reasons = [
    'USBHost_t36 library has ZERO USB Audio Class host driver. Checked the source.',
    'AudioInputUSB / AudioOutputUSB in the Teensy Audio Library are DEVICE mode only.',
    'USB Audio Class 2 uses isochronous transfers at 480 Mbps with clock feedback — far beyond USB MIDI bulk transfers.',
    'Even if the driver existed, handling 8×44.1kHz streams while running the sampler engine is borderline CPU-wise.',
    'The ESP32-S3 bridge in v2 was even worse: full-speed 12 Mbps USB only, 2-channel max experimental support.',
]
y = Inches(3.35)
for r in reasons:
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y + Inches(0.08),
                             Inches(0.1), Inches(0.1))
    dot.fill.solid(); dot.fill.fore_color.rgb = RED; dot.line.fill.background()
    add_text(s, Inches(0.95), y, Inches(11.9), Inches(0.4), r,
             size=11, color=WHITE, font='SF Pro Display')
    y += Inches(0.35)

add_text(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.4),
         'What the pros do instead',
         size=13, color=GREEN, bold=True, font='SF Pro Display')
add_text(s, Inches(0.5), Inches(5.85), Inches(12.3), Inches(1.2),
         'Dirtywave M8, Teenage Engineering OP-1, Polyend Tracker, Elektron Syntakt, Elektron '
         'Digitakt II — ALL of these professional portable instruments use an on-board codec, not '
         'USB audio hosting. When they need to integrate with a computer-based studio, they present '
         'themselves AS a USB audio device. That is the pattern Sampler-Crow follows: the Audio '
         'Shield SGTL5000 handles TRS I/O, and the Teensy\'s USB device port exposes us as an audio '
         'device to any Mac or PC alongside your Volt.',
         size=12, color=WHITE, font='SF Pro Display')

# ==========================================================
# SLIDE 14 — Roadmap
# ==========================================================
s = blank_slide()
add_title(s, 'Roadmap',
          'Where we are and where we are going')

phases = [
    ('✓', 'Phase 1-6',  'Foundation, sequencer, Launchpad integration',              GREEN),
    ('✓', 'Phase 7a',   '8-voice multi-timbral synth (kick/snare/hats/bass/lead)',    GREEN),
    ('✓', 'Phase 7b',   'Functional mixer: volume, mute, solo per track',             GREEN),
    ('→', 'Phase 7c',   'Velocity / pitch per step editing',                          ACCENT),
    ('→', 'Phase 8',    'Drum sampler: WAV playback from SD card (DrumMode port)',    ORANGE),
    ('  ', 'Phase 9',   'Pattern save/load to SD as JSON',                            DIM),
    ('  ', 'Phase 10',  'Mutable Instruments Plaits synth integration',               DIM),
    ('  ', 'Phase 11',  'Chord progression sequencer + scale quantizer',              DIM),
    ('  ', 'Phase 12',  'Port 5 sampler modes from new_hits (Pitch/Grain/Chop/Drum/Multi)', DIM),
    ('  ', 'Phase 13',  'CrowPanel 5" LVGL touchscreen UI',                           DIM),
    ('  ', 'Phase 14',  'Effects per track + AUX sends + master bus',                 DIM),
    ('  ', 'Phase 15',  'Live looping on audio tracks',                               DIM),
]
y = Inches(1.8)
for icon, phase, desc, color in phases:
    add_text(s, Inches(0.7), y, Inches(0.4), Inches(0.35),
             icon, size=16, color=color, bold=True, font='SF Pro Display')
    add_text(s, Inches(1.1), y, Inches(2.0), Inches(0.35),
             phase, size=13, color=color, bold=True, font='SF Pro Display')
    add_text(s, Inches(3.2), y, Inches(9.5), Inches(0.35),
             desc, size=12, color=WHITE, font='SF Pro Display')
    y += Inches(0.42)

# ==========================================================
# SLIDE 15 — Links & references
# ==========================================================
s = blank_slide()
add_title(s, 'Links & references',
          'Everything you need is on the repo')

add_text(s, Inches(0.5), Inches(1.8), Inches(12.3), Inches(0.4),
         'Project repository',
         size=14, color=ACCENT, bold=True, font='SF Pro Display')
add_hyperlink(s, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.4),
              'https://github.com/nohumans-com/Sampler-Crow',
              'https://github.com/nohumans-com/Sampler-Crow', size=13)

add_text(s, Inches(0.5), Inches(2.9), Inches(12.3), Inches(0.4),
         'Key documents in the repo',
         size=14, color=ACCENT, bold=True, font='SF Pro Display')
docs = [
    ('HARDWARE_DESIGN.md', 'Full v3 hardware spec — definitive source'),
    ('SOFTWARE_DESIGN.md', 'Original software architecture'),
    ('STATUS.md',          'Current build state and roadmap'),
    ('firmware/',          'Teensy 4.1 firmware (PlatformIO + Arduino)'),
    ('SamplerCrowApp/',    'macOS host app (SwiftUI + Swift 6)'),
]
y = Inches(3.3)
for name, desc in docs:
    add_text(s, Inches(0.7), y, Inches(3.5), Inches(0.35),
             name, size=12, color=YELLOW, font='SF Mono', bold=True)
    add_text(s, Inches(4.3), y, Inches(8.5), Inches(0.35),
             desc, size=12, color=WHITE, font='SF Pro Display')
    y += Inches(0.4)

add_text(s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.4),
         'Teensy 4.1 reference',
         size=14, color=ACCENT, bold=True, font='SF Pro Display')
refs = [
    ('PJRC Teensy 4.1',    'https://www.pjrc.com/store/teensy41.html'),
    ('Audio Library',      'https://www.pjrc.com/teensy/td_libs_Audio.html'),
    ('USBHost_t36',        'https://github.com/PaulStoffregen/USBHost_t36'),
    ('USB host cable',     'https://www.sparkfun.com/usb-host-cable-for-teensy-4-1-and-teensy-3-6.html'),
]
y = Inches(6.1)
for name, url in refs:
    add_text(s, Inches(0.7), y, Inches(3.5), Inches(0.3),
             name, size=11, color=WHITE, bold=True, font='SF Pro Display')
    add_hyperlink(s, Inches(4.3), y, Inches(8.5), Inches(0.3),
                  url, url, size=11)
    y += Inches(0.3)

# Save
import os
out = 'Sampler-Crow-Build-Guide.pptx'
prs.save(out)
print(f'Wrote {out}  ({os.path.getsize(out)//1024} KB, {len(prs.slides)} slides)')
